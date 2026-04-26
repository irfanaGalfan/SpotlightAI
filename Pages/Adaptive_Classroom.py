import streamlit as st
from PIL import Image
import os
import random
import firebase_admin
from firebase_admin import credentials, firestore
import docx2txt
from PyPDF2 import PdfReader
from pptx import Presentation
import time
from streamlit_lottie import st_lottie
import json

def show():
    # --- 1️⃣ Fallback Questions ---
    QUESTIONS = [
        {"id": 1, "question": "What year did the United Arab Emirates become an independent federation?",
         "options": ["1965", "1971", "1980", "1990"], "answer": "1971","difficulty":"Medium"},
        {"id": 2, "question": "Which city is home to the tallest building in the world, Burj Khalifa?",
         "options": ["Abu Dhabi", "Sharjah", "Dubai", "Ajman"], "answer": "Dubai","difficulty":"Easy"},
        {"id": 3, "question": "Which of these is a traditional Emirati sport involving trained birds?",
         "options": ["Falconry", "Camel Racing", "Horse Racing", "Dhow Sailing"], "answer": "Falconry","difficulty":"Medium"},
        {"id": 4, "question": "What is the official currency of the United Arab Emirates?",
         "options": ["Riyal", "Dirham", "Dinar", "Pound"], "answer": "Dirham","difficulty":"Easy"},
        {"id": 5, "question": "Which desert covers much of the UAE’s land area?",
         "options": ["Sahara Desert", "Gobi Desert", "Rub' al Khali", "Kalahari Desert"], "answer": "Rub' al Khali", "difficulty":"Easy"},
        {"id": 6, "question": "Which ancient archaeological site in the UAE shows human burial practices dating back thousands of years?",
         "options": ["Al Jahili Fort", "Jebel Buhais", "Al Ain Oasis", "Sheikh Zayed Grand Mosque"], "answer": "Jebel Buhais","difficulty":"Hard"},
        {"id": 7, "question": "The national day of the UAE is celebrated on which date?",
         "options": ["January 1", "June 30", "December 2", "October 15"], "answer": "December 2","difficulty":"Medium"},
        {"id": 8, "question": "Which animal’s milk was a key part of the traditional nomadic diet in the UAE?",
         "options": ["Camel milk", "Goat milk", "Cow milk", "Sheep milk"], "answer": "Camel milk","difficulty":"Hard"},
        {"id": 9, "question": "Which emirate is the capital of the UAE?",
         "options": ["Dubai", "Sharjah", "Abu Dhabi", "Ras Al Khaimah"], "answer": "Abu Dhabi","difficulty":"Easy"},
        {"id": 10, "question": "Environmental sustainability in the UAE is exemplified by which planned eco-city?",
         "options": ["Green City Abu Dhabi", "Masdar City", "Desert City Dubai", "Eco Oasis Sharjah"], "answer": "Masdar City","difficulty":"Hard"},
    ]

    # --- 2️⃣ Firebase initialization ---
    @st.cache_resource
    def init_firebase():
        if not firebase_admin._apps:
            cred = credentials.Certificate(
                "C:/Users/Administrator/OneDrive/Desktop/AI_Project/aiproject-4094d-firebase-adminsdk-fbsvc-c46223dc00.json"
            )
            firebase_admin.initialize_app(cred)
        return firestore.client()

    db = init_firebase()

    try:
        list(db.collections())
        st.success("✅ Firebase connected successfully")
    except Exception as e:
        st.error(f"❌ Firebase connection failed: {e}")
        st.stop()

    # --- 3️⃣ Session state initialization ---
    for key, default in [('quiz_active', False), ('current_index', 0),
                         ('quiz_students', []), ('asked_questions', set()),
                         ('current_question', {}), ('material_text', "")]:
        if key not in st.session_state:
            st.session_state[key] = default
    # --- RL Q-table initialization ---
    # --- Q_table structure: ---
    # {
    #   "LOW":    {"Easy": 0.0, "Medium": 0.0, "Hard": 0.0},
    #   "MEDIUM": {"Easy": 0.0, "Medium": 0.0, "Hard": 0.0},
    #   "HIGH":   {"Easy": 0.0, "Medium": 0.0, "Hard": 0.0}
    # }

    if "Q_table" not in st.session_state:
        st.session_state["Q_table"] = {}   # {state: {action: value}}

    ACTIONS = ["Easy", "Medium", "Hard"]

    def get_student_state(student):
        if student["score"] <= 1:
            return "LOW"
        elif student["score"] <= 3:
            return "MEDIUM"
        else:
            return "HIGH"
    def rl_select_question(student, questions, epsilon=0.2):
        state = get_student_state(student)
        Q = st.session_state["Q_table"]

        if state not in Q:
            Q[state] = {a: 0.0 for a in ACTIONS}

        # --- ε-greedy selection ---
        if random.random() < epsilon:
            chosen_difficulty = random.choice(ACTIONS)
        else:
            chosen_difficulty = max(Q[state], key=Q[state].get)

        # --- RULE: prevent difficulty jump ---
        if state == "LOW" and chosen_difficulty == "Hard":
            chosen_difficulty = "Easy"

        # --- Select question ---
        filtered = [
            q for q in questions
            if q["difficulty"] == chosen_difficulty
            and q["question"] not in st.session_state["asked_questions"]
        ]

        if filtered:
            return random.choice(filtered), state, chosen_difficulty
        else:
            return random.choice(questions), state, chosen_difficulty


    # --- Folder containing student images ---
    IMAGE_FOLDER = r"C:\Users\Administrator\OneDrive\Desktop\AI_Project\classes"

    st.title("AI-Based Adaptive Quiz System")

    # --- 4️⃣ Select Class ---
    #class_docs = db.collection("classes").stream()
    #classes = [{"id": doc.id, "name": doc.to_dict().get("class_name")} for doc in class_docs]

    #selected_class = st.selectbox(
     #   "Select Class",
      #  options=classes,
       # format_func=lambda x: x["name"]
    #)
    class_docs = db.collection("classes").stream()
    classes = []

    for doc in class_docs:
        classes.append({
            "id": doc.id,
            "name": doc.id   # ✅ document ID IS the class name
        })

    if not classes:
        st.error("❌ No classes found. Please create a class first.")
        st.stop()

    selected_class = st.selectbox(
        "Select Class",
        options=classes,
        format_func=lambda x: x["name"]
    )

    # --- 5️⃣ Fetch students ---
    def fetch_students_from_class(class_name):
        students = []

        students_ref = (
            db.collection("classes")
              .document(class_name)
              .collection("students")
        )

        for doc in students_ref.stream():
            data = doc.to_dict()
            students.append({
                "name": doc.id,
                "score": data.get("score", 0),
                "turns": data.get("turns", 0)
            })

        return students

    if selected_class:
        students = fetch_students_from_class(selected_class["name"])
        st.subheader("Students in Class")
        st.table(students)

    # --- 6️⃣ Upload learning material ---
    uploaded_file = st.file_uploader("Upload material (PDF / PPT / DOCX)", type=['pdf','pptx','docx'])
    def extract_text(file):
        text = ""
        if uploaded_file.name.endswith(".pdf"):
            pdf = PdfReader(file)
            for page in pdf.pages:
                text += page.extract_text()
        elif uploaded_file.name.endswith(".docx"):
            text = docx2txt.process(file)
        elif uploaded_file.name.endswith(".pptx"):
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        return text
    if uploaded_file:
        st.session_state['material_text'] = extract_text(uploaded_file)
        st.success("Material text extracted successfully")

    # --- 7️⃣ Slider for number of students ---
    if selected_class and students and len(students) > 0:
        max_students = min(10, len(students))
        num_students = st.slider(
            "Select number of students for this session",
            min_value=1,
            max_value=max_students,
            value=1)
    else:
        st.warning("⚠️ No students found in this class. Add students first.")
        st.stop()

    #num_students = st.slider("Select number of students for this session", 1, min(10, len(students)), 3)

    # --- 8️⃣ Generate session ---
    if st.button("Generate Session"):
        # Select students with minimum turns
        sorted_students = sorted(students, key=lambda s: s['turns'])
        st.session_state['quiz_students'] = sorted_students[:num_students]
        st.session_state['current_index'] = 0
        st.session_state['asked_questions'] = set()
        st.session_state['current_question'] = {}
        st.session_state['quiz_active'] = True
        st.success("Session generated! Starting quiz...")
        audio_url = "https://www.soundhelix.com/examples/mp3/SoundHelix-Song-1.mp3"
        audio_html = f"""
                <audio autoplay>
                  <source src="{audio_url}" type="audio/mpeg">
                </audio>
                """
        st.components.v1.html(audio_html, height=0, width=0)

    # --- 9️⃣ Question selection ---
    #def dqn_select_question(student_name, questions, material_text=None):
        #for q in questions:
            #if q["question"] not in st.session_state['asked_questions']:
                #return q
        #return random.choice(questions)

    # --- 10️⃣ Quiz container ---
    quiz_container = st.container()

    def display_next_student():
          
        if not st.session_state.get('quiz_active', False):
            return

        current_index = st.session_state['current_index']
        if current_index >= len(st.session_state['quiz_students']):
            st.success("🏆 Quiz completed!")
            st.session_state['quiz_active'] = False
         # --- Stop hidden audio ---
            st.components.v1.html("""
            <script>
                const audio = document.getElementById("quiz_audio");
                if (audio) {
                    audio.pause();
                    audio.currentTime = 0;
                }
            </script>
            """, height=0, width=0)
            return
        

        student = st.session_state['quiz_students'][current_index]

        # Pick or reuse question
        if student['name'] in st.session_state['current_question']:
            question = st.session_state['current_question'][student['name']]
        else:
            #question = dqn_select_question(student['name'], QUESTIONS, st.session_state['material_text'])
            question, rl_state, rl_action = rl_select_question(student, QUESTIONS)
            st.session_state["rl_context"] = (rl_state, rl_action)
            st.session_state['current_question'][student['name']] = question
            st.session_state['asked_questions'].add(question['question'])

        with quiz_container:
            # --- Student image ---
            img_path = None
            class_folder = os.path.join(IMAGE_FOLDER, selected_class["name"])
            for ext in ["jpg", "png", "jpeg"]:
                possible_path = os.path.join(class_folder, f"{student['name']}.{ext}")
                if os.path.exists(possible_path):
                    img_path = possible_path
                    break
            if img_path:
                img = Image.open(img_path)
                col1, col2, col3 = st.columns([1, 2, 1])
                with col2:
                    st.image(img, width=150)
            else:
                st.text("No photo available")

            # --- Student info & question ---
            st.subheader(f"👤 Student: {student['name']}")
            st.write(f"Score: {student['score']} | Turns: {student['turns']}")
            st.write(f"❓ Question: {question['question']}")
            selected_answer = st.radio("Select your answer:", question['options'], key=f"answer_{current_index}")

            if st.button("Submit", key=f"submit_{current_index}"):
                correct = selected_answer == question['answer']
                # --- Show feedback notification ---
                if correct:
                    st.success("✅ Correct answer!")
                else:
                    st.error(f"❌ Wrong answer! Correct answer: {question['answer']}")

                # --- 🔹 RL Reward & Q-Table Update ---
                rl_state, rl_action = st.session_state.get("rl_context", (None, None))
                if rl_state is not None and rl_action is not None:
                    # Simple reward
                    reward = 10 if correct else -5

                    # Optional: confidence / difficulty rule
                    if rl_state == "LOW" and rl_action == "Easy" and correct:
                        reward += 5

                    # Q-learning update
                    alpha = 0.3
                    Q = st.session_state["Q_table"]
                    if rl_state not in Q:
                        Q[rl_state] = {a: 0.0 for a in ["Easy", "Medium", "Hard"]}
                    Q[rl_state][rl_action] += alpha * (reward - Q[rl_state][rl_action])
                
                    # --- 🔹 DEBUG: Show Q-table ---
                    st.write("🔹 Q-table:", st.session_state["Q_table"])
                # --- Update Firebase ---
                # Inside Submit button, before updating Firebase
                student_ref = (
                                    db.collection("classes")
                                      .document(selected_class["name"])
                                      .collection("students")
                                      .document(student["name"])
                                )

                student_doc = student_ref.get()
                if student_doc.exists:
                    s_data = student_doc.to_dict()
                    answered_questions = s_data.get("answered_questions", [])
                    recent_scores = s_data.get("recent_scores", [])
                    cumulative_turns = s_data.get("turns", 0)
                    cumulative_score = s_data.get("score", 0)
                    difficulty_history = s_data.get("difficulty_history", [])  # <-- fetch it
                else:
                    answered_questions = []
                    recent_scores = []
                    cumulative_turns = 0
                    cumulative_score = 0
                    difficulty_history = []  # initialize if missing
                # Then append new answer and update score as you currently do
                # --- 🔹 Append new answer and update score here ---
                answered_questions.append(question['question'])
                recent_scores.append(1 if correct else -1)
                cumulative_turns += 1
                cumulative_score = sum(1 for s in recent_scores if s == 1)
                difficulty_history.append(question['difficulty'])
                st.write("🔹 Difficulty History before saving:", difficulty_history)

                student_ref.set({
                    "answered_questions": answered_questions,
                    "recent_scores": recent_scores,
                    "score": cumulative_score,
                    "turns": cumulative_turns,
                    "difficulty_history": difficulty_history  # <-- save it
                }, merge=True)

                # --- Update session variables ---
                student['score'] = cumulative_score
                student['turns'] = cumulative_turns
                del st.session_state['current_question'][student['name']]
                st.session_state['current_index'] += 1
                # --- Clear container and display next student ---
                quiz_container.empty()
                display_next_student()

    # --- 11️⃣ Start quiz if active ---
    if st.session_state.get('quiz_active', False):
        display_next_student()
